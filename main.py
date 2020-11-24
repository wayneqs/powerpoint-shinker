from pptx import Presentation
from pptx.shapes.picture import Picture
from PIL import Image
import io

def to_inch(emu):
  emu_to_inch = 1/914400.0
  return round(emu * emu_to_inch, 2)

def calc_target_size(shape):
  dpi, _ = shape.image.dpi
  return (int(to_inch(shape.width) * dpi), int(to_inch(shape.height) * dpi))

prs = Presentation('presso.pptx')
for slide in prs.slides: 
  for shape in slide.shapes:
    if type(shape) is Picture:
      print(shape.image.size)
      with Image.open(io.BytesIO(shape.image.blob)) as im:
        im = im.convert('RGB').resize(calc_target_size(shape))
        im.save('map.jpg', 'JPEG')

# prs.save("presso_BOOM.pptx")